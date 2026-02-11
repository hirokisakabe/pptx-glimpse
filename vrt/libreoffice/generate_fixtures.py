#!/usr/bin/env python3
"""
LibreOffice VRT 用の PPTX フィクスチャを python-pptx で生成する。

Usage:
    python3 vrt/libreoffice/generate_fixtures.py
"""

import os

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


def make_element(tag, **attribs):
    """OOXML 要素を作成する (e.g., make_element("a:gradFill"))"""
    element = etree.SubElement(etree.Element("dummy"), qn(tag))
    element.getparent().remove(element)
    for key, val in attribs.items():
        element.set(key, val)
    return element

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "fixtures")

SLIDE_WIDTH = 9144000
SLIDE_HEIGHT = 5143500

# 4:3 slide size
SLIDE_WIDTH_4_3 = 9144000
SLIDE_HEIGHT_4_3 = 6858000


def new_presentation():
    """スライドサイズを固定した新しいプレゼンテーションを生成する。"""
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH)
    prs.slide_height = Emu(SLIDE_HEIGHT)
    return prs


def new_presentation_4_3():
    """4:3 スライドサイズの新しいプレゼンテーションを生成する。"""
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH_4_3)
    prs.slide_height = Emu(SLIDE_HEIGHT_4_3)
    return prs


def create_basic_shapes():
    """基本図形: rect, ellipse, roundRect (ソリッド塗り + 枠線)"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    shapes_def = [
        (MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(2.5), Inches(2),
         RGBColor(0x44, 0x72, 0xC4), "Rectangle"),
        (MSO_SHAPE.OVAL, Inches(3.5), Inches(0.5), Inches(2.5), Inches(2),
         RGBColor(0xED, 0x7D, 0x31), "Oval"),
        (MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(0.5), Inches(2.5), Inches(2),
         RGBColor(0xA5, 0xA5, 0xA5), "Rounded Rect"),
        (MSO_SHAPE.DIAMOND, Inches(0.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0xFF, 0xC0, 0x00), "Diamond"),
        (MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0x5B, 0x9B, 0xD5), "Triangle"),
        (MSO_SHAPE.HEXAGON, Inches(6.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0x70, 0xAD, 0x47), "Hexagon"),
    ]

    for shape_type, left, top, width, height, color, _label in shapes_def:
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    prs.save(os.path.join(OUTPUT_DIR, "lo-basic-shapes.pptx"))
    print("  Created: lo-basic-shapes.pptx")


def create_text_formatting():
    """テキスト書式: 太字、イタリック、フォントサイズ、色、配置"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tests = [
        {"text": "Bold Text", "bold": True, "size": Pt(24)},
        {"text": "Italic Text", "italic": True, "size": Pt(24)},
        {"text": "Large 36pt", "size": Pt(36)},
        {"text": "Red Text", "size": Pt(24), "color": RGBColor(0xFF, 0x00, 0x00)},
        {"text": "Center Aligned", "size": Pt(24), "align": PP_ALIGN.CENTER},
        {"text": "Right Aligned", "size": Pt(24), "align": PP_ALIGN.RIGHT},
    ]

    for i, t in enumerate(tests):
        col = i % 2
        row = i // 2
        left = Inches(0.3 + col * 4.7)
        top = Inches(0.3 + row * 1.7)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(4.4), Inches(1.4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        shape.line.width = Pt(0.5)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t["text"]

        if "align" in t:
            p.alignment = t["align"]

        run = p.runs[0]
        run.font.name = "Liberation Sans"
        run.font.size = t.get("size", Pt(18))

        if t.get("bold"):
            run.font.bold = True
        if t.get("italic"):
            run.font.italic = True
        if "color" in t:
            run.font.color.rgb = t["color"]

    prs.save(os.path.join(OUTPUT_DIR, "lo-text-formatting.pptx"))
    print("  Created: lo-text-formatting.pptx")


def create_fill_and_lines():
    """塗りと線: ソリッド塗り各色 + 枠線スタイル"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    colors = [
        (RGBColor(0xFF, 0x63, 0x84), "Pink"),
        (RGBColor(0x36, 0xA2, 0xEB), "Blue"),
        (RGBColor(0xFF, 0xCE, 0x56), "Yellow"),
    ]

    for i, (color, _label) in enumerate(colors):
        left = Inches(0.5 + i * 3)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(0.5),
            Inches(2.5), Inches(2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(2)

    # 枠線のみの図形（塗りなし）
    no_fill_shapes = [
        (MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0x44, 0x72, 0xC4), Pt(1)),
        (MSO_SHAPE.OVAL, Inches(3.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0xED, 0x7D, 0x31), Pt(2)),
        (MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0x70, 0xAD, 0x47), Pt(3)),
    ]

    for shape_type, left, top, width, height, line_color, line_width in no_fill_shapes:
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.background()
        shape.line.color.rgb = line_color
        shape.line.width = line_width

    prs.save(os.path.join(OUTPUT_DIR, "lo-fill-and-lines.pptx"))
    print("  Created: lo-fill-and-lines.pptx")


def create_gradient_fills():
    """グラデーション塗り: 水平/垂直/対角/3色/枠線付き/暗色"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    gradients = [
        {
            "shape": MSO_SHAPE.ROUNDED_RECTANGLE,
            "stops": [(0, "4472C4"), (100000, "ED7D31")],
            "angle": 0,
            "label": "Horizontal",
        },
        {
            "shape": MSO_SHAPE.ROUNDED_RECTANGLE,
            "stops": [(0, "FF0000"), (100000, "FFFF00")],
            "angle": 5400000,
            "label": "Vertical",
        },
        {
            "shape": MSO_SHAPE.ROUNDED_RECTANGLE,
            "stops": [(0, "70AD47"), (100000, "7030A0")],
            "angle": 2700000,
            "label": "Diagonal",
        },
        {
            "shape": MSO_SHAPE.RECTANGLE,
            "stops": [(0, "FF0000"), (50000, "FFFFFF"), (100000, "4472C4")],
            "angle": 0,
            "label": "3-Color",
        },
        {
            "shape": MSO_SHAPE.OVAL,
            "stops": [(0, "5B9BD5"), (100000, "FFC000")],
            "angle": 0,
            "label": "With Border",
            "border": True,
        },
        {
            "shape": MSO_SHAPE.DIAMOND,
            "stops": [(0, "000000"), (100000, "808080")],
            "angle": 5400000,
            "label": "Dark",
        },
    ]

    for i, g in enumerate(gradients):
        col = i % 3
        row = i // 3
        left = Inches(0.3 + col * 3.2)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(
            g["shape"], left, top, Inches(2.8), Inches(2.2)
        )
        shape.fill.background()

        # XML でグラデーション塗りを設定
        spPr = shape._element.spPr
        # noFill を削除
        for child in list(spPr):
            tag = etree.QName(child.tag).localname
            if tag in ("noFill", "solidFill"):
                spPr.remove(child)

        gradFill = make_element("a:gradFill")
        gsLst = make_element("a:gsLst")
        for pos, color_hex in g["stops"]:
            gs = make_element("a:gs")
            gs.set("pos", str(pos))
            srgbClr = make_element("a:srgbClr")
            srgbClr.set("val", color_hex)
            gs.append(srgbClr)
            gsLst.append(gs)
        gradFill.append(gsLst)

        lin = make_element("a:lin")
        lin.set("ang", str(g["angle"]))
        lin.set("scaled", "1")
        gradFill.append(lin)
        spPr.append(gradFill)

        if g.get("border"):
            shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
            shape.line.width = Pt(2)

    prs.save(os.path.join(OUTPUT_DIR, "lo-gradient-fills.pptx"))
    print("  Created: lo-gradient-fills.pptx")


def create_dash_lines():
    """線のダッシュスタイル: solid, dash, dot, dashDot, lgDash, lgDashDot"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    dash_styles = [
        ("solid", "Solid"),
        ("dash", "Dash"),
        ("dot", "Dot"),
        ("dashDot", "DashDot"),
        ("lgDash", "LgDash"),
        ("lgDashDot", "LgDashDot"),
    ]

    for i, (dash_val, label) in enumerate(dash_styles):
        col = i % 3
        row = i // 3
        left = Inches(0.3 + col * 3.2)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(2.8), Inches(2.2)
        )
        shape.fill.background()
        shape.line.color.rgb = RGBColor(0x44, 0x72, 0xC4)
        shape.line.width = Pt(2.5)

        # XML でダッシュスタイルを設定
        ln = shape.line._ln
        prstDash = make_element("a:prstDash")
        prstDash.set("val", dash_val)
        ln.append(prstDash)

        # ラベルテキスト
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Liberation Sans"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(os.path.join(OUTPUT_DIR, "lo-dash-lines.pptx"))
    print("  Created: lo-dash-lines.pptx")


def create_text_decoration():
    """テキスト装飾: 下線、取り消し線、太字+下線、イタリック+取り消し線、複数行、上付き"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tests = [
        {"text": "Underline Text", "underline": True},
        {"text": "Strikethrough", "strikethrough": True},
        {"text": "Bold + Underline", "bold": True, "underline": True},
        {"text": "Italic + Strike", "italic": True, "strikethrough": True},
        {"text": "Line 1\nLine 2\nLine 3", "multiline": True},
        {"text": "Normal", "superscript": "sup", "sup_text": "super"},
    ]

    for i, t in enumerate(tests):
        col = i % 3
        row = i // 3
        left = Inches(0.3 + col * 3.2)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(2.8), Inches(2.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        shape.line.width = Pt(0.5)

        tf = shape.text_frame
        tf.word_wrap = True

        if t.get("multiline"):
            lines = t["text"].split("\n")
            for j, line_text in enumerate(lines):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line_text
                run = p.runs[0]
                run.font.name = "Liberation Sans"
                run.font.size = Pt(18)
        elif t.get("superscript"):
            p = tf.paragraphs[0]
            p.text = t["text"]
            run = p.runs[0]
            run.font.name = "Liberation Sans"
            run.font.size = Pt(20)
            # 上付き文字を追加
            sup_run = p.add_run()
            sup_run.text = t["sup_text"]
            sup_run.font.name = "Liberation Sans"
            sup_run.font.size = Pt(20)
            rPr = sup_run._r.get_or_add_rPr()
            rPr.set("baseline", "30000")
        else:
            p = tf.paragraphs[0]
            p.text = t["text"]
            run = p.runs[0]
            run.font.name = "Liberation Sans"
            run.font.size = Pt(20)

            if t.get("bold"):
                run.font.bold = True
            if t.get("italic"):
                run.font.italic = True
            if t.get("underline"):
                run.font.underline = True
            if t.get("strikethrough"):
                rPr = run._r.get_or_add_rPr()
                rPr.set("strike", "sngStrike")

    prs.save(os.path.join(OUTPUT_DIR, "lo-text-decoration.pptx"))
    print("  Created: lo-text-decoration.pptx")


def create_tables():
    """テーブル: ヘッダ行 + データ行 (交互背景色)"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    rows, cols = 4, 3
    left = Inches(1)
    top = Inches(0.8)
    width = Inches(8)
    height = Inches(3.5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    headers = ["Name", "Category", "Value"]
    data = [
        ["Alpha", "Type A", "100"],
        ["Beta", "Type B", "200"],
        ["Gamma", "Type C", "300"],
    ]

    # ヘッダ行
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Liberation Sans"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # データ行
    for i, row_data in enumerate(data):
        bg = RGBColor(0xF2, 0xF2, 0xF2) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        for j, value in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0]
            run.font.name = "Liberation Sans"
            run.font.size = Pt(14)

    prs.save(os.path.join(OUTPUT_DIR, "lo-tables.pptx"))
    print("  Created: lo-tables.pptx")


def create_bullets():
    """箇条書き・番号付きリスト: buChar (丸/ダッシュ) + buAutoNum (数字/アルファベット)"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bullet_configs = [
        {
            "title": "Bullet (dot)",
            "type": "buChar",
            "char": "\u2022",
            "items": ["First item", "Second item", "Third item"],
        },
        {
            "title": "Bullet (dash)",
            "type": "buChar",
            "char": "-",
            "items": ["Apple", "Banana", "Cherry"],
        },
        {
            "title": "Numbered (1. 2. 3.)",
            "type": "buAutoNum",
            "scheme": "arabicPeriod",
            "items": ["Step one", "Step two", "Step three"],
        },
        {
            "title": "Alpha (a. b. c.)",
            "type": "buAutoNum",
            "scheme": "alphaLcPeriod",
            "items": ["Option A", "Option B", "Option C"],
        },
    ]

    for i, config in enumerate(bullet_configs):
        col = i % 2
        row = i // 2
        left = Inches(0.3 + col * 4.7)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(4.4), Inches(2.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        shape.line.width = Pt(0.5)

        tf = shape.text_frame
        tf.word_wrap = True

        # タイトル行
        p = tf.paragraphs[0]
        p.text = config["title"]
        run = p.runs[0]
        run.font.name = "Liberation Sans"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)

        # 箇条書き項目
        for item_text in config["items"]:
            p = tf.add_paragraph()
            p.text = item_text
            p.level = 0
            run = p.runs[0]
            run.font.name = "Liberation Sans"
            run.font.size = Pt(14)

            pPr = p._pPr
            if pPr is None:
                pPr = make_element("a:pPr")
                p._p.insert(0, pPr)
            pPr.set("marL", str(Inches(0.3)))
            pPr.set("indent", str(-Inches(0.2)))

            if config["type"] == "buChar":
                buChar = make_element("a:buChar")
                buChar.set("char", config["char"])
                pPr.append(buChar)
            elif config["type"] == "buAutoNum":
                buAutoNum = make_element("a:buAutoNum")
                buAutoNum.set("type", config["scheme"])
                pPr.append(buAutoNum)

    prs.save(os.path.join(OUTPUT_DIR, "lo-bullets.pptx"))
    print("  Created: lo-bullets.pptx")


def create_transforms():
    """回転・フリップ: 45度/90度回転、水平/垂直フリップ、組み合わせ"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    transforms = [
        {"rotation": 45.0, "flipH": False, "flipV": False, "label": "Rot 45"},
        {"rotation": 90.0, "flipH": False, "flipV": False, "label": "Rot 90"},
        {"rotation": 0.0, "flipH": True, "flipV": False, "label": "FlipH"},
        {"rotation": 0.0, "flipH": False, "flipV": True, "label": "FlipV"},
        {"rotation": 0.0, "flipH": True, "flipV": True, "label": "FlipHV"},
        {"rotation": 30.0, "flipH": True, "flipV": False, "label": "Rot30+FlipH"},
    ]

    for i, t in enumerate(transforms):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 3.2)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, left, top, Inches(2.5), Inches(1.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1)

        if t["rotation"] != 0.0:
            shape.rotation = t["rotation"]

        xfrm = shape._element.spPr.xfrm
        if t["flipH"]:
            xfrm.set("flipH", "1")
        if t["flipV"]:
            xfrm.set("flipV", "1")

    prs.save(os.path.join(OUTPUT_DIR, "lo-transforms.pptx"))
    print("  Created: lo-transforms.pptx")


def create_groups():
    """グループ図形: 2グループ (各2図形)"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # python-pptx でのグループ作成は XML 操作が必要
    # グループ1: 青矩形 + オレンジ楕円
    group_shapes = [
        {
            "group_left": Inches(0.5),
            "group_top": Inches(0.8),
            "group_w": Inches(4),
            "group_h": Inches(3.5),
            "children": [
                {
                    "type": MSO_SHAPE.RECTANGLE,
                    "left": Inches(0.5),
                    "top": Inches(0.8),
                    "w": Inches(1.8),
                    "h": Inches(1.5),
                    "color": RGBColor(0x44, 0x72, 0xC4),
                },
                {
                    "type": MSO_SHAPE.OVAL,
                    "left": Inches(2.5),
                    "top": Inches(2.5),
                    "w": Inches(1.8),
                    "h": Inches(1.5),
                    "color": RGBColor(0xED, 0x7D, 0x31),
                },
            ],
        },
        {
            "group_left": Inches(5.2),
            "group_top": Inches(0.8),
            "group_w": Inches(4),
            "group_h": Inches(3.5),
            "children": [
                {
                    "type": MSO_SHAPE.DIAMOND,
                    "left": Inches(5.2),
                    "top": Inches(0.8),
                    "w": Inches(1.8),
                    "h": Inches(1.5),
                    "color": RGBColor(0x70, 0xAD, 0x47),
                },
                {
                    "type": MSO_SHAPE.ISOSCELES_TRIANGLE,
                    "left": Inches(7.2),
                    "top": Inches(2.5),
                    "w": Inches(1.8),
                    "h": Inches(1.5),
                    "color": RGBColor(0xFF, 0xC0, 0x00),
                },
            ],
        },
    ]

    for group_def in group_shapes:
        # グループ図形を XML で構築
        spTree = slide.shapes._spTree

        grpSp = make_element("p:grpSp")

        # grpSpPr (グループの位置・サイズ)
        grpSpPr = make_element("p:grpSpPr")
        xfrm = make_element("a:xfrm")
        off = make_element("a:off")
        off.set("x", str(int(group_def["group_left"])))
        off.set("y", str(int(group_def["group_top"])))
        xfrm.append(off)
        ext = make_element("a:ext")
        ext.set("cx", str(int(group_def["group_w"])))
        ext.set("cy", str(int(group_def["group_h"])))
        xfrm.append(ext)
        chOff = make_element("a:chOff")
        chOff.set("x", str(int(group_def["group_left"])))
        chOff.set("y", str(int(group_def["group_top"])))
        xfrm.append(chOff)
        chExt = make_element("a:chExt")
        chExt.set("cx", str(int(group_def["group_w"])))
        chExt.set("cy", str(int(group_def["group_h"])))
        xfrm.append(chExt)
        grpSpPr.append(xfrm)
        grpSp.append(grpSpPr)

        # nvGrpSpPr
        nvGrpSpPr = make_element("p:nvGrpSpPr")
        cNvPr = make_element("p:cNvPr")
        cNvPr.set("id", str(100 + group_shapes.index(group_def)))
        cNvPr.set("name", f"Group {group_shapes.index(group_def) + 1}")
        nvGrpSpPr.append(cNvPr)
        cNvGrpSpPr = make_element("p:cNvGrpSpPr")
        nvGrpSpPr.append(cNvGrpSpPr)
        nvPr = make_element("p:nvPr")
        nvGrpSpPr.append(nvPr)
        grpSp.insert(0, nvGrpSpPr)

        # 子図形を追加
        for ci, child in enumerate(group_def["children"]):
            sp = make_element("p:sp")

            # nvSpPr
            nvSpPr = make_element("p:nvSpPr")
            cNvPr2 = make_element("p:cNvPr")
            cNvPr2.set("id", str(200 + group_shapes.index(group_def) * 10 + ci))
            cNvPr2.set("name", f"Shape {ci + 1}")
            nvSpPr.append(cNvPr2)
            cNvSpPr = make_element("p:cNvSpPr")
            nvSpPr.append(cNvSpPr)
            nvPr2 = make_element("p:nvPr")
            nvSpPr.append(nvPr2)
            sp.append(nvSpPr)

            # spPr
            spPr = make_element("p:spPr")
            childXfrm = make_element("a:xfrm")
            childOff = make_element("a:off")
            childOff.set("x", str(int(child["left"])))
            childOff.set("y", str(int(child["top"])))
            childXfrm.append(childOff)
            childExt = make_element("a:ext")
            childExt.set("cx", str(int(child["w"])))
            childExt.set("cy", str(int(child["h"])))
            childXfrm.append(childExt)
            spPr.append(childXfrm)

            # prstGeom (プリセット図形)
            prstGeom = make_element("a:prstGeom")
            shape_map = {
                MSO_SHAPE.RECTANGLE: "rect",
                MSO_SHAPE.OVAL: "ellipse",
                MSO_SHAPE.DIAMOND: "diamond",
                MSO_SHAPE.ISOSCELES_TRIANGLE: "triangle",
            }
            prstGeom.set("prst", shape_map[child["type"]])
            avLst = make_element("a:avLst")
            prstGeom.append(avLst)
            spPr.append(prstGeom)

            # solidFill
            solidFill = make_element("a:solidFill")
            srgbClr = make_element("a:srgbClr")
            color = child["color"]
            srgbClr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
            solidFill.append(srgbClr)
            spPr.append(solidFill)

            # ln (枠線)
            ln = make_element("a:ln")
            ln.set("w", str(int(Pt(1.5))))
            lnFill = make_element("a:solidFill")
            lnClr = make_element("a:srgbClr")
            lnClr.set("val", "333333")
            lnFill.append(lnClr)
            ln.append(lnFill)
            spPr.append(ln)

            sp.append(spPr)
            grpSp.append(sp)

        spTree.append(grpSp)

    prs.save(os.path.join(OUTPUT_DIR, "lo-groups.pptx"))
    print("  Created: lo-groups.pptx")


def create_slide_background():
    """スライド背景: ソリッドカラー背景 + 白テキストボックス + 枠線のみ図形"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # スライド背景色を設定
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)

    # 白い矩形 + テキスト
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(3.5), Inches(1.5)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape1.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape1.line.width = Pt(0)
    tf = shape1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "White Box"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Liberation Sans"
    run.font.size = Pt(24)
    run.font.bold = True

    # 黄色い矩形
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1), Inches(3.5), Inches(1.5)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0xFF, 0xCE, 0x56)
    shape2.line.width = Pt(0)
    tf2 = shape2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Yellow Box"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Liberation Sans"
    run2.font.size = Pt(24)
    run2.font.bold = True

    # 枠線のみの図形 (白枠)
    shape3 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(2.5), Inches(3), Inches(5), Inches(1.8)
    )
    shape3.fill.background()
    shape3.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape3.line.width = Pt(3)

    prs.save(os.path.join(OUTPUT_DIR, "lo-slide-background.pptx"))
    print("  Created: lo-slide-background.pptx")


def create_flowchart_shapes():
    """フローチャート図形: 8つのフローチャートプリセット (ソリッド塗り + 枠線)"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    flowchart_defs = [
        (MSO_SHAPE.FLOWCHART_PROCESS, RGBColor(0x44, 0x72, 0xC4), "Process"),
        (MSO_SHAPE.FLOWCHART_ALTERNATE_PROCESS, RGBColor(0xED, 0x7D, 0x31), "Alt Process"),
        (MSO_SHAPE.FLOWCHART_DECISION, RGBColor(0xA5, 0xA5, 0xA5), "Decision"),
        (MSO_SHAPE.FLOWCHART_DATA, RGBColor(0xFF, 0xC0, 0x00), "Data"),
        (MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS, RGBColor(0x5B, 0x9B, 0xD5), "Predefined"),
        (MSO_SHAPE.FLOWCHART_DOCUMENT, RGBColor(0x70, 0xAD, 0x47), "Document"),
        (MSO_SHAPE.FLOWCHART_TERMINATOR, RGBColor(0xFF, 0x63, 0x84), "Terminator"),
        (MSO_SHAPE.FLOWCHART_PREPARATION, RGBColor(0x36, 0xA2, 0xEB), "Preparation"),
    ]

    for i, (shape_type, color, _label) in enumerate(flowchart_defs):
        col = i % 4
        row = i // 4
        left = Inches(0.3 + col * 2.4)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(shape_type, left, top, Inches(2.1), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    prs.save(os.path.join(OUTPUT_DIR, "lo-flowchart-shapes.pptx"))
    print("  Created: lo-flowchart-shapes.pptx")


def create_arrows_stars():
    """矢印と星形図形: MSO_SHAPE presets"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    arrow_defs = [
        (MSO_SHAPE.LEFT_RIGHT_ARROW, RGBColor(0x44, 0x72, 0xC4)),
        (MSO_SHAPE.UP_DOWN_ARROW, RGBColor(0xED, 0x7D, 0x31)),
        (MSO_SHAPE.QUAD_ARROW, RGBColor(0xA5, 0xA5, 0xA5)),
        (MSO_SHAPE.CIRCULAR_ARROW, RGBColor(0xFF, 0xC0, 0x00)),
    ]

    for i, (shape_type, color) in enumerate(arrow_defs):
        left = Inches(0.3 + i * 2.4)
        shape = slide.shapes.add_shape(shape_type, left, Inches(0.3), Inches(2.1), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    star_defs = [
        (MSO_SHAPE.STAR_6_POINT, RGBColor(0xFF, 0x63, 0x84)),
        (MSO_SHAPE.STAR_8_POINT, RGBColor(0x36, 0xA2, 0xEB)),
        (MSO_SHAPE.STAR_10_POINT, RGBColor(0xFF, 0xCE, 0x56)),
        (MSO_SHAPE.STAR_12_POINT, RGBColor(0x99, 0x66, 0xFF)),
        (MSO_SHAPE.STAR_16_POINT, RGBColor(0x4B, 0xC0, 0xC0)),
    ]

    for i, (shape_type, color) in enumerate(star_defs):
        left = Inches(0.3 + i * 1.9)
        shape = slide.shapes.add_shape(shape_type, left, Inches(2.8), Inches(1.6), Inches(1.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    prs.save(os.path.join(OUTPUT_DIR, "lo-arrows-stars.pptx"))
    print("  Created: lo-arrows-stars.pptx")


def create_callouts_arcs():
    """吹き出しとアーク図形"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    callout_defs = [
        (MSO_SHAPE.RECTANGULAR_CALLOUT, RGBColor(0x44, 0x72, 0xC4)),
        (MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT, RGBColor(0xED, 0x7D, 0x31)),
        (MSO_SHAPE.OVAL_CALLOUT, RGBColor(0xA5, 0xA5, 0xA5)),
        (MSO_SHAPE.CLOUD_CALLOUT, RGBColor(0xFF, 0xC0, 0x00)),
    ]

    for i, (shape_type, color) in enumerate(callout_defs):
        col = i % 2
        row = i // 2
        left = Inches(0.3 + col * 4.7)
        top = Inches(0.3 + row * 1.8)
        shape = slide.shapes.add_shape(shape_type, left, top, Inches(4.4), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    arc_defs = [
        (MSO_SHAPE.ARC, RGBColor(0x5B, 0x9B, 0xD5)),
        (MSO_SHAPE.BLOCK_ARC, RGBColor(0x70, 0xAD, 0x47)),
        (MSO_SHAPE.CHORD, RGBColor(0xFF, 0x63, 0x84)),
    ]

    for i, (shape_type, color) in enumerate(arc_defs):
        left = Inches(0.3 + i * 3.2)
        shape = slide.shapes.add_shape(shape_type, left, Inches(3.9), Inches(2.8), Inches(1.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    prs.save(os.path.join(OUTPUT_DIR, "lo-callouts-arcs.pptx"))
    print("  Created: lo-callouts-arcs.pptx")


def create_math_other():
    """数式記号とその他の図形"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    math_defs = [
        (MSO_SHAPE.MATH_PLUS, RGBColor(0x44, 0x72, 0xC4)),
        (MSO_SHAPE.MATH_MINUS, RGBColor(0xED, 0x7D, 0x31)),
        (MSO_SHAPE.MATH_MULTIPLY, RGBColor(0xA5, 0xA5, 0xA5)),
        (MSO_SHAPE.MATH_DIVIDE, RGBColor(0xFF, 0xC0, 0x00)),
        (MSO_SHAPE.MATH_EQUAL, RGBColor(0x5B, 0x9B, 0xD5)),
        (MSO_SHAPE.MATH_NOT_EQUAL, RGBColor(0x70, 0xAD, 0x47)),
    ]

    for i, (shape_type, color) in enumerate(math_defs):
        left = Inches(0.3 + i * 1.55)
        shape = slide.shapes.add_shape(shape_type, left, Inches(0.3), Inches(1.3), Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    other_defs = [
        (MSO_SHAPE.GEAR_6, RGBColor(0xFF, 0x63, 0x84)),
        (MSO_SHAPE.GEAR_9, RGBColor(0x36, 0xA2, 0xEB)),
        (MSO_SHAPE.DONUT, RGBColor(0xFF, 0xCE, 0x56)),
        (MSO_SHAPE.NO_SYMBOL, RGBColor(0x99, 0x66, 0xFF)),
        (MSO_SHAPE.CROSS, RGBColor(0x4B, 0xC0, 0xC0)),
        (MSO_SHAPE.HEART, RGBColor(0xF0, 0x80, 0x80)),
    ]

    for i, (shape_type, color) in enumerate(other_defs):
        left = Inches(0.3 + i * 1.55)
        shape = slide.shapes.add_shape(shape_type, left, Inches(1.9), Inches(1.3), Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    prs.save(os.path.join(OUTPUT_DIR, "lo-math-other.pptx"))
    print("  Created: lo-math-other.pptx")


def create_image():
    """画像埋め込み: Pillow で動的生成した画像を配置"""
    import shutil
    import tempfile

    from PIL import Image, ImageDraw

    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    temp_dir = tempfile.mkdtemp()

    try:
        # 画像1: カラーグラデーション
        img1 = Image.new("RGB", (200, 200))
        pixels = img1.load()
        for y in range(200):
            for x in range(200):
                pixels[x, y] = (int(255 * x / 200), int(255 * y / 200), 128)
        img1_path = os.path.join(temp_dir, "gradient.png")
        img1.save(img1_path)

        # 画像2: チェッカーボード
        img2 = Image.new("RGB", (200, 200), "white")
        draw2 = ImageDraw.Draw(img2)
        for row in range(10):
            for col in range(10):
                if (row + col) % 2 == 0:
                    x1, y1 = col * 20, row * 20
                    draw2.rectangle([x1, y1, x1 + 20, y1 + 20], fill=(0, 0, 0))
        img2_path = os.path.join(temp_dir, "checker.png")
        img2.save(img2_path)

        # 画像3: 円パターン
        img3 = Image.new("RGB", (200, 200), "white")
        draw3 = ImageDraw.Draw(img3)
        draw3.ellipse([10, 10, 190, 190], fill=(68, 114, 196), outline=(51, 51, 51))
        draw3.ellipse([60, 60, 140, 140], fill=(237, 125, 49))
        img3_path = os.path.join(temp_dir, "circles.png")
        img3.save(img3_path)

        slide.shapes.add_picture(img1_path, Inches(0.5), Inches(0.5), Inches(2.8), Inches(4))
        slide.shapes.add_picture(img2_path, Inches(3.6), Inches(0.5), Inches(2.8), Inches(4))
        slide.shapes.add_picture(img3_path, Inches(6.7), Inches(0.5), Inches(2.8), Inches(4))

        prs.save(os.path.join(OUTPUT_DIR, "lo-image.pptx"))
        print("  Created: lo-image.pptx")
    finally:
        shutil.rmtree(temp_dir)


def create_charts():
    """チャート: 縦棒、折れ線、円、横棒"""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 縦棒グラフ
    chart_data1 = CategoryChartData()
    chart_data1.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data1.add_series("Sales", (19.2, 21.4, 16.7, 20.8))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.3), Inches(0.3), Inches(4.3), Inches(2.2),
        chart_data1,
    )

    # 折れ線グラフ
    chart_data2 = CategoryChartData()
    chart_data2.categories = ["Jan", "Feb", "Mar", "Apr"]
    chart_data2.add_series("Revenue", (30, 35, 32, 40))
    chart_data2.add_series("Profit", (10, 12, 11, 15))
    slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(5.0), Inches(0.3), Inches(4.3), Inches(2.2),
        chart_data2,
    )

    # 円グラフ
    chart_data3 = CategoryChartData()
    chart_data3.categories = ["Red", "Green", "Blue"]
    chart_data3.add_series("Share", (40, 35, 25))
    slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(0.3), Inches(2.8), Inches(4.3), Inches(2.2),
        chart_data3,
    )

    # 横棒グラフ
    chart_data4 = CategoryChartData()
    chart_data4.categories = ["Alpha", "Beta", "Gamma"]
    chart_data4.add_series("Values", (50, 65, 45))
    slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(5.0), Inches(2.8), Inches(4.3), Inches(2.2),
        chart_data4,
    )

    prs.save(os.path.join(OUTPUT_DIR, "lo-charts.pptx"))
    print("  Created: lo-charts.pptx")


def create_connectors():
    """コネクタ線: XML 直接操作で p:cxnSp を生成"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    spTree = slide.shapes._spTree

    connectors_def = [
        {
            "id": "10",
            "name": "Straight",
            "prst": "line",
            "x": "500000",
            "y": "500000",
            "cx": "3000000",
            "cy": "0",
            "color": "4472C4",
            "width": "25400",
        },
        {
            "id": "11",
            "name": "Dashed",
            "prst": "line",
            "x": "500000",
            "y": "1500000",
            "cx": "3000000",
            "cy": "2000000",
            "color": "ED7D31",
            "width": "25400",
            "dash": "dash",
        },
        {
            "id": "12",
            "name": "Dotted",
            "prst": "line",
            "x": "5000000",
            "y": "500000",
            "cx": "0",
            "cy": "4000000",
            "color": "70AD47",
            "width": "25400",
            "dash": "dot",
        },
    ]

    for c in connectors_def:
        cxnSp = make_element("p:cxnSp")

        nvCxnSpPr = make_element("p:nvCxnSpPr")
        cNvPr = make_element("p:cNvPr")
        cNvPr.set("id", c["id"])
        cNvPr.set("name", c["name"])
        nvCxnSpPr.append(cNvPr)
        nvCxnSpPr.append(make_element("p:cNvCxnSpPr"))
        nvCxnSpPr.append(make_element("p:nvPr"))
        cxnSp.append(nvCxnSpPr)

        spPr = make_element("p:spPr")

        xfrm = make_element("a:xfrm")
        off = make_element("a:off")
        off.set("x", c["x"])
        off.set("y", c["y"])
        xfrm.append(off)
        ext = make_element("a:ext")
        ext.set("cx", c["cx"])
        ext.set("cy", c["cy"])
        xfrm.append(ext)
        spPr.append(xfrm)

        prstGeom = make_element("a:prstGeom")
        prstGeom.set("prst", c["prst"])
        prstGeom.append(make_element("a:avLst"))
        spPr.append(prstGeom)

        ln = make_element("a:ln")
        ln.set("w", c["width"])
        solidFill = make_element("a:solidFill")
        srgbClr = make_element("a:srgbClr")
        srgbClr.set("val", c["color"])
        solidFill.append(srgbClr)
        ln.append(solidFill)
        if "dash" in c:
            prstDash = make_element("a:prstDash")
            prstDash.set("val", c["dash"])
            ln.append(prstDash)
        spPr.append(ln)

        cxnSp.append(spPr)
        spTree.append(cxnSp)

    prs.save(os.path.join(OUTPUT_DIR, "lo-connectors.pptx"))
    print("  Created: lo-connectors.pptx")


def _build_custom_path(path_elem, commands):
    """カスタムジオメトリのパスコマンドを構築するヘルパー"""
    for cmd in commands:
        if cmd[0] == "moveTo":
            moveTo = make_element("a:moveTo")
            pt = make_element("a:pt")
            pt.set("x", str(cmd[1]))
            pt.set("y", str(cmd[2]))
            moveTo.append(pt)
            path_elem.append(moveTo)
        elif cmd[0] == "lnTo":
            lnTo = make_element("a:lnTo")
            pt = make_element("a:pt")
            pt.set("x", str(cmd[1]))
            pt.set("y", str(cmd[2]))
            lnTo.append(pt)
            path_elem.append(lnTo)
        elif cmd[0] == "cubicBezTo":
            bez = make_element("a:cubicBezTo")
            for j in range(1, 4):
                pt = make_element("a:pt")
                pt.set("x", str(cmd[j * 2 - 1]))
                pt.set("y", str(cmd[j * 2]))
                bez.append(pt)
            path_elem.append(bez)
        elif cmd[0] == "close":
            path_elem.append(make_element("a:close"))


def create_custom_geometry():
    """カスタムジオメトリ: a:custGeom でカスタムパスを定義"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    shapes_def = [
        {
            "left": Inches(0.5),
            "top": Inches(0.5),
            "w": Inches(3.5),
            "h": Inches(3.5),
            "color": "FFC000",
            "commands": [
                ("moveTo", 500, 0),
                ("lnTo", 650, 350),
                ("lnTo", 1000, 400),
                ("lnTo", 750, 650),
                ("lnTo", 800, 1000),
                ("lnTo", 500, 850),
                ("lnTo", 200, 1000),
                ("lnTo", 250, 650),
                ("lnTo", 0, 400),
                ("lnTo", 350, 350),
                ("close",),
            ],
        },
        {
            "left": Inches(5.5),
            "top": Inches(0.5),
            "w": Inches(3.5),
            "h": Inches(3.5),
            "color": "5B9BD5",
            "commands": [
                ("moveTo", 0, 500),
                ("cubicBezTo", 250, 0, 750, 1000, 1000, 500),
                ("close",),
            ],
        },
    ]

    spTree = slide.shapes._spTree

    for idx, s in enumerate(shapes_def):
        sp = make_element("p:sp")

        nvSpPr = make_element("p:nvSpPr")
        cNvPr = make_element("p:cNvPr")
        cNvPr.set("id", str(100 + idx))
        cNvPr.set("name", f"Custom {idx + 1}")
        nvSpPr.append(cNvPr)
        nvSpPr.append(make_element("p:cNvSpPr"))
        nvSpPr.append(make_element("p:nvPr"))
        sp.append(nvSpPr)

        spPr = make_element("p:spPr")

        xfrm = make_element("a:xfrm")
        off = make_element("a:off")
        off.set("x", str(int(s["left"])))
        off.set("y", str(int(s["top"])))
        xfrm.append(off)
        ext = make_element("a:ext")
        ext.set("cx", str(int(s["w"])))
        ext.set("cy", str(int(s["h"])))
        xfrm.append(ext)
        spPr.append(xfrm)

        custGeom = make_element("a:custGeom")
        custGeom.append(make_element("a:avLst"))
        custGeom.append(make_element("a:gdLst"))
        custGeom.append(make_element("a:ahLst"))
        custGeom.append(make_element("a:cxnLst"))
        rect = make_element("a:rect")
        rect.set("l", "0")
        rect.set("t", "0")
        rect.set("r", "0")
        rect.set("b", "0")
        custGeom.append(rect)

        pathLst = make_element("a:pathLst")
        path = make_element("a:path")
        path.set("w", "1000")
        path.set("h", "1000")
        _build_custom_path(path, s["commands"])
        pathLst.append(path)
        custGeom.append(pathLst)
        spPr.append(custGeom)

        solidFill = make_element("a:solidFill")
        srgbClr = make_element("a:srgbClr")
        srgbClr.set("val", s["color"])
        solidFill.append(srgbClr)
        spPr.append(solidFill)

        ln = make_element("a:ln")
        ln.set("w", "12700")
        lnFill = make_element("a:solidFill")
        lnClr = make_element("a:srgbClr")
        lnClr.set("val", "333333")
        lnFill.append(lnClr)
        ln.append(lnFill)
        spPr.append(ln)

        sp.append(spPr)
        spTree.append(sp)

    prs.save(os.path.join(OUTPUT_DIR, "lo-custom-geometry.pptx"))
    print("  Created: lo-custom-geometry.pptx")


def create_slide_size_4_3():
    """4:3 スライドサイズ: 基本図形 + テキスト + 背景色を 1 スライドに配置"""
    prs = new_presentation_4_3()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # スライド背景色
    background = slide.background
    fill_bg = background.fill
    fill_bg.solid()
    fill_bg.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)

    # 上段: 基本図形 3 つ
    shapes_def = [
        (MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.3), Inches(2.8), Inches(2.5),
         RGBColor(0x44, 0x72, 0xC4)),
        (MSO_SHAPE.OVAL, Inches(3.5), Inches(0.3), Inches(2.8), Inches(2.5),
         RGBColor(0xED, 0x7D, 0x31)),
        (MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(0.3), Inches(2.8), Inches(2.5),
         RGBColor(0x70, 0xAD, 0x47)),
    ]

    for shape_type, left, top, width, height, color in shapes_def:
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    # 中段: テキストボックス 2 つ
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(3.1), Inches(4.5), Inches(2.2)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
    left_box.line.color.rgb = RGBColor(0x44, 0x72, 0xC4)
    left_box.line.width = Pt(1)
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "4:3 Layout Text"
    p_l.alignment = PP_ALIGN.CENTER
    run_l = p_l.runs[0]
    run_l.font.name = "Liberation Sans"
    run_l.font.size = Pt(20)
    run_l.font.bold = True
    run_l.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(3.1), Inches(4.5), Inches(2.2)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)
    right_box.line.color.rgb = RGBColor(0xFF, 0xC0, 0x00)
    right_box.line.width = Pt(1)
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "Right Content"
    p_r.alignment = PP_ALIGN.CENTER
    run_r = p_r.runs[0]
    run_r.font.name = "Liberation Sans"
    run_r.font.size = Pt(20)
    run_r.font.italic = True
    run_r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # 下段: フッターバー
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(5.6), Inches(9.2), Inches(1.6)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(0x44, 0x54, 0x6A)
    footer.line.width = Pt(0)
    tf_f = footer.text_frame
    tf_f.word_wrap = True
    p_f = tf_f.paragraphs[0]
    p_f.text = "Footer on 4:3 slide"
    p_f.alignment = PP_ALIGN.CENTER
    run_f = p_f.runs[0]
    run_f.font.name = "Liberation Sans"
    run_f.font.size = Pt(14)
    run_f.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save(os.path.join(OUTPUT_DIR, "lo-slide-size-4-3.pptx"))
    print("  Created: lo-slide-size-4-3.pptx")


def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    print("Generating LibreOffice VRT fixtures...")
    create_basic_shapes()
    create_text_formatting()
    create_fill_and_lines()
    create_gradient_fills()
    create_dash_lines()
    create_text_decoration()
    create_tables()
    create_bullets()
    create_transforms()
    create_groups()
    create_slide_background()
    create_flowchart_shapes()
    create_arrows_stars()
    create_callouts_arcs()
    create_math_other()
    create_image()
    create_charts()
    create_connectors()
    create_custom_geometry()
    create_slide_size_4_3()
    print("Done!")


if __name__ == "__main__":
    main()

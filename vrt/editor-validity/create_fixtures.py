#!/usr/bin/env python3
"""
Generate PPTX fixtures for editor-validity tests with python-pptx.

The fixtures are source / expected pairs consumed by editor-validity.test.ts,
plus a basic-shapes source used by slide topology and shape add/delete checks.
basic-shapes.pptx duplicates the renderer VRT fixture of the same name on
purpose: the editor-validity suite must stay runnable without generating the
vrt/libreoffice/ fixture set.

Usage:
    python3 vrt/editor-validity/create_fixtures.py
"""

import base64
import os
import re
import tempfile
import zipfile

from pptx import Presentation
from pptx.chart.data import BubbleChartData, CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from PIL import Image

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "fixtures")

SLIDE_WIDTH = 9144000
SLIDE_HEIGHT = 5143500


def new_presentation():
    """Generate a new presentation with fixed slide sizes."""
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH)
    prs.slide_height = Emu(SLIDE_HEIGHT)
    return prs


def create_basic_shapes():
    """Basic shapes: 6 presets from rect to hexagon (solid fill + border)"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    shapes_def = [
        (MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(2.5), Inches(2),
         RGBColor(0x44, 0x72, 0xC4), "Rectangle"),
        (MSO_SHAPE.OVAL, Inches(3.5), Inches(0.5), Inches(2.5), Inches(2),
         RGBColor(0xED, 0x7D, 0x31), "Oval"),
        (MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(0.5), Inches(2.5), Inches(2),
         RGBColor(0xA5, 0xA5, 0xA5), "Rounded Rect"),
        (MSO_SHAPE.DIAMOND, Inches(0.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0xFF, 0xC0, 0x00), "Diamond"),
        (MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0x5B, 0x9B, 0xD5), "Triangle"),
        (MSO_SHAPE.HEXAGON, Inches(6.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0x70, 0xAD, 0x47), "Hexagon"),
    ]

    for shape_type, left, top, width, height, color, _label in shapes_def:
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    prs.save(os.path.join(OUTPUT_DIR, "basic-shapes.pptx"))
    print("  Created: basic-shapes.pptx")


def create_editor_validity_text_fixture(filename, text):
    """Fixture pair for editor text replacement validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.5))
    title.text_frame.text = "LibreOffice editor validity: text"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.5), Inches(7.8), Inches(1.8)
    )
    shape.name = "Editable Text Target"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE2, 0xF0, 0xD9)
    shape.line.color.rgb = RGBColor(0x70, 0xAD, 0x47)
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    run = paragraph.runs[0]
    run.font.name = "Liberation Sans"
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_transform_fixture(filename, left, top, width, height):
    """Fixture pair for editor move / resize validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.5))
    title.text_frame.text = "LibreOffice editor validity: transform"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.name = "Move Resize Target"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x5B, 0x9B, 0xD5)
    shape.line.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    shape.line.width = Pt(2)

    tf = shape.text_frame
    paragraph = tf.paragraphs[0]
    paragraph.text = "Move + resize"
    run = paragraph.runs[0]
    run.font.name = "Liberation Sans"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    marker = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(7.7), Inches(3.9), Inches(0.45), Inches(0.45)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = RGBColor(0xED, 0x7D, 0x31)
    marker.line.fill.background()

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_formatting_fixture(filename, *, expected):
    """Fixture pair for editor run property validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.5))
    title.text_frame.text = "LibreOffice editor validity: formatting"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(8.2), Inches(2.0)
    )
    shape.name = "Editable Formatting Target"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)
    shape.line.color.rgb = RGBColor(0xBF, 0x90, 0x00)
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    paragraph = tf.paragraphs[0]
    paragraph.text = "Editable formatting target"
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0]

    if expected:
        run.font.name = "Liberation Serif"
        run.font.size = Pt(30)
        run.font.bold = False
        run.font.italic = False
        run.font.underline = False
        run.font.color.rgb = RGBColor(0x9C, 0x00, 0x00)
    else:
        run.font.name = "Liberation Sans"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.italic = True
        run.font.underline = True
        run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_paragraph_fixture(filename, *, expected):
    """Fixture pair for editor paragraph property validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.5))
    title.text_frame.text = "LibreOffice editor validity: paragraph"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.3), Inches(8.0), Inches(2.0)
    )
    shape.name = "Editable Paragraph Target"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xDE, 0xEB, 0xF7)
    shape.line.color.rgb = RGBColor(0x2F, 0x75, 0xB5)
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    paragraph = tf.paragraphs[0]
    paragraph.text = "Paragraph properties target"
    paragraph.alignment = PP_ALIGN.RIGHT if expected else PP_ALIGN.LEFT
    paragraph.level = 1 if expected else 0
    if expected:
        p_pr = paragraph._p.get_or_add_pPr()
        bu_char = OxmlElement("a:buChar")
        bu_char.set("char", "\u2022")
        p_pr.insert(0, bu_char)
    run = paragraph.runs[0]
    run.font.name = "Liberation Sans"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_image_fixture(filename, image_base64):
    """Fixture pair for editor image replacement validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.5))
    title.text_frame.text = "LibreOffice editor validity: image"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        tmp.write(base64.b64decode(image_base64))
        image_path = tmp.name

    try:
        pic = slide.shapes.add_picture(
            image_path, Inches(2.2), Inches(1.4), Inches(4.8), Inches(2.7)
        )
        pic.name = "Replace Image Target"
    finally:
        os.unlink(image_path)

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_picture_crop_fixture(filename, *, expected):
    """Fixture pair for existing stretch picture crop validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.5))
    title.text_frame.text = "LibreOffice editor validity: picture crop"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        image = Image.new("RGB", (120, 80))
        for y in range(80):
            for x in range(120):
                image.putpixel(
                    (x, y),
                    (
                        240 if x < 60 else 30,
                        220 if y < 40 else 40,
                        80 if x < 60 else 220,
                    ),
                )
        image.save(tmp, format="PNG")
        image_path = tmp.name

    try:
        pic = slide.shapes.add_picture(
            image_path, Inches(1.5), Inches(1.2), Inches(7.0), Inches(3.6)
        )
        pic.name = "Picture Crop Target"
        if expected:
            pic.crop_left = 0.25
            pic.crop_top = 0.10
            pic.crop_right = 0.05
            pic.crop_bottom = 0.15
    finally:
        os.unlink(image_path)

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_chart_fixture(filename, *, expected):
    """Fixture pair for existing chart data update validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["Apr", "May", "Jun"] if expected else ["Jan", "Feb"]
    if expected:
        chart_data.add_series("Edited revenue", (40, 55, 70))
        chart_data.add_series("Edited cost", (25, 30, 42))
        chart_data.add_series("Edited profit", (15, 25, 28))
    else:
        chart_data.add_series("Revenue", (10, 20))
        chart_data.add_series("Cost", (7, 12))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(0.7),
        Inches(8.4),
        Inches(4.5),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "LibreOffice editor validity: chart"
    chart.has_legend = True
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "Month"
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Amount"

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_chart_removal_fixture(filename, *, expected):
    """Fixture pair for existing chart series removal validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["Apr", "May", "Jun"] if expected else ["Jan", "Feb"]
    if expected:
        chart_data.add_series("Edited revenue", (40, 55, 70))
    else:
        chart_data.add_series("Revenue", (10, 20))
        chart_data.add_series("Cost", (7, 12))
        chart_data.add_series("Profit", (3, 8))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(0.7),
        Inches(8.4),
        Inches(4.5),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "LibreOffice editor validity: chart removal"
    chart.has_legend = True
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "Month"
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Amount"

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_category_combo_chart_fixture(filename, *, expected):
    """Fixture pair for fixed-topology bar + line category combo updates."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["Apr", "May", "Jun"] if expected else ["Jan", "Feb"]
    if expected:
        chart_data.add_series("Edited columns", (40, 55, 70))
        chart_data.add_series("Edited trend", (35, 48, 63))
    else:
        chart_data.add_series("Columns", (10, 20))
        chart_data.add_series("Trend", (7, 12))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(0.7),
        Inches(8.4),
        Inches(4.5),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "LibreOffice editor validity: category combo chart"
    chart.has_legend = True

    plot_area = chart._element.find(qn("c:chart")).find(qn("c:plotArea"))
    bar_chart = plot_area.find(qn("c:barChart"))
    series = bar_chart.findall(qn("c:ser"))
    if len(series) != 2:
        raise RuntimeError("combo fixture requires two category series")
    line_series = series[1]
    bar_chart.remove(line_series)
    line_chart = OxmlElement("c:lineChart")
    grouping = OxmlElement("c:grouping")
    grouping.set("val", "standard")
    line_chart.append(grouping)
    vary_colors = OxmlElement("c:varyColors")
    vary_colors.set("val", "0")
    line_chart.append(vary_colors)
    line_chart.append(line_series)
    for axis_id in bar_chart.findall(qn("c:axId")):
        copied_axis_id = OxmlElement("c:axId")
        copied_axis_id.set("val", axis_id.get("val"))
        line_chart.append(copied_axis_id)
    plot_area.insert(plot_area.index(bar_chart) + 1, line_chart)

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_scatter_chart_fixture(filename, *, expected):
    """Fixture pair for existing scatter chart XY data update validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = XyChartData()
    series_values = (
        [
            ("Edited revenue", [(1.0, 40.0), (2.0, 55.0), (3.0, 70.0)]),
            ("Edited cost", [(10.0, 25.0), (20.0, 42.0)]),
        ]
        if expected
        else [
            ("Revenue", [(1.0, 10.0), (2.0, 20.0)]),
            ("Cost", [(3.0, 7.0), (4.0, 12.0)]),
        ]
    )
    for name, points in series_values:
        series = chart_data.add_series(name)
        for x_value, y_value in points:
            series.add_data_point(x_value, y_value)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(0.8),
        Inches(0.7),
        Inches(8.4),
        Inches(4.5),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "LibreOffice editor validity: scatter chart"
    chart.has_legend = True
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "X"
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Y"

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_bubble_chart_fixture(filename, *, expected):
    """Fixture pair for existing bubble chart XYZ data update validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = BubbleChartData()
    series_values = (
        [
            ("Edited revenue", [(1.0, 40.0, 5.0), (2.0, 55.0, 8.0), (3.0, 70.0, 13.0)]),
            ("Edited cost", [(10.0, 25.0, 21.0), (20.0, 42.0, 34.0)]),
        ]
        if expected
        else [
            ("Revenue", [(1.0, 10.0, 4.0), (2.0, 20.0, 8.0)]),
            ("Cost", [(3.0, 7.0, 6.0), (4.0, 12.0, 9.0)]),
        ]
    )
    for name, points in series_values:
        series = chart_data.add_series(name)
        for x_value, y_value, size in points:
            series.add_data_point(x_value, y_value, size)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BUBBLE,
        Inches(0.8),
        Inches(0.7),
        Inches(8.4),
        Inches(4.5),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "LibreOffice editor validity: bubble chart"
    chart.has_legend = True
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.has_title = True
    chart.category_axis.axis_title.text_frame.text = "X"
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Y"

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_table_text_fixture(filename, text):
    """Fixture pair for existing Table cell text replacement validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.5))
    title.text_frame.text = "LibreOffice editor validity: Table cell text"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    table = slide.shapes.add_table(
        2,
        2,
        Inches(1.0),
        Inches(1.3),
        Inches(8.0),
        Inches(2.2),
    ).table
    values = [[text, "Unedited sibling"], ["Preserved row", "Preserved cell"]]
    for row_index, row in enumerate(table.rows):
        for column_index, cell in enumerate(row.cells):
            cell.text = values[row_index][column_index]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            run = paragraph.runs[0]
            run.font.name = "Liberation Sans"
            run.font.size = Pt(22)
            run.font.bold = row_index == 0
            run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def set_table_cell_border(cell, side, width, color):
    """Set one inline DrawingML Table cell border."""
    tc_pr = cell._tc.get_or_add_tcPr()
    tag = f"a:ln{side}"
    for child in list(tc_pr):
        if child.tag.endswith(f"ln{side}"):
            tc_pr.remove(child)
    line = OxmlElement(tag)
    line.set("w", str(width))
    solid_fill = OxmlElement("a:solidFill")
    srgb = OxmlElement("a:srgbClr")
    srgb.set("val", color)
    solid_fill.append(srgb)
    line.append(solid_fill)
    tc_pr.insert(0, line)


def create_editor_validity_table_cell_properties_fixture(filename, *, expected):
    """Fixture pair for existing Table cell fill/border/margin validity checks."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.5))
    title.text_frame.text = "LibreOffice editor validity: Table cell properties"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    table = slide.shapes.add_table(
        2,
        2,
        Inches(1.0),
        Inches(1.3),
        Inches(8.0),
        Inches(2.4),
    ).table
    values = [["Edited property target", "Unedited sibling"], ["Preserved row", "Preserved cell"]]
    for row_index, row in enumerate(table.rows):
        for column_index, cell in enumerate(row.cells):
            cell.text = values[row_index][column_index]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            run = paragraph.runs[0]
            run.font.name = "Liberation Sans"
            run.font.size = Pt(20)

    target = table.cell(0, 0)
    target.fill.solid()
    target.fill.fore_color.rgb = RGBColor(0xF4, 0xB1, 0x83) if expected else RGBColor(0xD9, 0xEA, 0xF7)
    target.margin_left = Emu(457200 if expected else 91440)
    target.margin_right = Emu(91440)
    target.margin_top = Emu(182880 if expected else 91440)
    target.margin_bottom = Emu(91440)
    set_table_cell_border(target, "L", 25400 if expected else 12700, "C00000" if expected else "4472C4")

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_group_fixture(filename, *, grouped):
    """Fixture pair for lossless existing group / ungroup topology edits."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.5))
    title.text_frame.text = "LibreOffice editor validity: group topology"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    first = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.4), Inches(3.0), Inches(1.4)
    )
    first.name = "Group Target 1"
    first.fill.solid()
    first.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
    first.line.color.rgb = RGBColor(0x20, 0x38, 0x64)
    first.line.width = Pt(1.5)

    second = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(5.0), Inches(2.5), Inches(3.0), Inches(1.4)
    )
    second.name = "Group Target 2"
    second.fill.solid()
    second.fill.fore_color.rgb = RGBColor(0xED, 0x7D, 0x31)
    second.line.color.rgb = RGBColor(0x84, 0x36, 0x10)
    second.line.width = Pt(1.5)

    if grouped:
        group_left = Inches(1.0)
        group_top = Inches(1.4)
        group_width = Inches(7.0)
        group_height = Inches(2.5)
        group = OxmlElement("p:grpSp")
        non_visual = OxmlElement("p:nvGrpSpPr")
        properties = OxmlElement("p:cNvPr")
        properties.set("id", "100")
        properties.set("name", "Expected Group")
        non_visual.append(properties)
        non_visual.append(OxmlElement("p:cNvGrpSpPr"))
        non_visual.append(OxmlElement("p:nvPr"))
        group.append(non_visual)

        group_properties = OxmlElement("p:grpSpPr")
        transform = OxmlElement("a:xfrm")
        for tag, attributes in [
            ("a:off", {"x": group_left, "y": group_top}),
            ("a:ext", {"cx": group_width, "cy": group_height}),
            ("a:chOff", {"x": group_left, "y": group_top}),
            ("a:chExt", {"cx": group_width, "cy": group_height}),
        ]:
            element = OxmlElement(tag)
            for name, value in attributes.items():
                element.set(name, str(int(value)))
            transform.append(element)
        group_properties.append(transform)
        group.append(group_properties)
        group.append(first._element)
        group.append(second._element)
        slide.shapes._spTree.append(group)

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_affine_move_fixture():
    """Native group with rotation, flip, and non-uniform child mapping for move validity."""
    source_path = os.path.join(OUTPUT_DIR, "editor-validity-group-expected.pptx")
    output_path = os.path.join(OUTPUT_DIR, "editor-validity-affine-move.pptx")
    with zipfile.ZipFile(source_path, "r") as source_archive:
        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as output_archive:
            for info in source_archive.infolist():
                data = source_archive.read(info.filename)
                if info.filename == "ppt/slides/slide1.xml":
                    xml = data.decode("utf-8")
                    group_match = re.search(
                        r'<p:grpSp\b[^>]*>[\s\S]*?name="Expected Group"[\s\S]*?</p:grpSp>',
                        xml,
                    )
                    if group_match is None:
                        raise RuntimeError("affine move group XML was not found")
                    group_xml = group_match.group(0)
                    xfrm_match = re.search(r"<a:xfrm[^>]*>[\s\S]*?</a:xfrm>", group_xml)
                    if xfrm_match is None:
                        raise RuntimeError("affine move group transform XML was not found")
                    xfrm = xfrm_match.group(0)
                    ext_match = re.search(r'<a:ext cx="(\d+)" cy="(\d+)"/>', xfrm)
                    if ext_match is None:
                        raise RuntimeError("affine move group extents were not found")
                    width = int(ext_match.group(1))
                    transformed = re.sub(
                        r"<a:xfrm[^>]*>",
                        '<a:xfrm rot="5400000" flipH="1">',
                        xfrm,
                        count=1,
                    ).replace(
                        ext_match.group(0),
                        f'<a:ext cx="{width * 2}" cy="{ext_match.group(2)}"/>',
                    )
                    xml = xml.replace(group_xml, group_xml.replace(xfrm, transformed))
                    data = xml.encode("utf-8")
                output_archive.writestr(info, data)
    print("  Created: editor-validity-affine-move.pptx")


def create_editor_validity_cross_slide_move_fixture(filename, *, expected):
    """Two-slide pair for destination-local drawing identity and writer validity."""
    prs = new_presentation()
    source_slide = prs.slides.add_slide(prs.slide_layouts[6])
    destination_slide = prs.slides.add_slide(prs.slide_layouts[6])

    anchor = destination_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.8), Inches(1.0), Inches(3.0), Inches(3.8)
    )
    anchor.name = "Destination Anchor"
    anchor.fill.solid()
    anchor.fill.fore_color.rgb = RGBColor(0xED, 0x7D, 0x31)

    moved_owner = destination_slide if expected else source_slide
    moved = moved_owner.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(3.2), Inches(1.8)
    )
    moved.name = "Cross Slide Move Target"
    moved.fill.solid()
    moved.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
    moved.text_frame.text = "Moved between slides"
    moved.text_frame.paragraphs[0].runs[0].font.size = Pt(18)

    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("Values", (2, 5, 3))
    moved_chart_frame = moved_owner.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(3.5),
        Inches(3.2),
        Inches(2.1),
        chart_data,
    )
    moved_chart_frame.name = "Cross Slide Chart Move Target"
    moved_chart = moved_chart_frame.chart
    moved_chart.has_title = True
    moved_chart.chart_title.text_frame.text = "Moved chart"
    moved_chart.has_legend = False

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_drawing_delete_fixture(filename):
    """Fixture containing one native picture, Table, Chart, and group delete target."""
    grouped_path = os.path.join(OUTPUT_DIR, "editor-validity-group-expected.pptx")
    prs = Presentation(grouped_path)
    slide = prs.slides[0]

    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        tmp.write(base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAQAAAAECAIAAAAmkwkpAAAAEUlEQVR4nGP8z4AATEhsPBwAM9EBBzDn4UwAAAAASUVORK5CYII="
        ))
        image_path = tmp.name
    try:
        picture = slide.shapes.add_picture(
            image_path, Inches(0.5), Inches(4.1), Inches(1.3), Inches(0.8)
        )
        picture.name = "Delete Picture"
    finally:
        os.unlink(image_path)

    table = slide.shapes.add_table(
        1, 1, Inches(2.0), Inches(4.1), Inches(2.0), Inches(0.8)
    )
    table.name = "Delete Table"
    table.table.cell(0, 0).text = "Delete"

    chart_data = CategoryChartData()
    chart_data.categories = ["A"]
    chart_data.add_series("Delete", (1,))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(4.3),
        Inches(3.8),
        Inches(2.0),
        Inches(1.2),
        chart_data,
    )
    chart.name = "Delete Chart"

    prs.save(os.path.join(OUTPUT_DIR, filename))
    print(f"  Created: {filename}")


def create_editor_validity_theme_fixture(filename, *, expected):
    """Fixture pair for existing theme color- and font-scheme edits."""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(7.6), Inches(2.6)
    )
    shape.name = "Theme Accent Target"
    shape.fill.solid()
    shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    shape.line.color.theme_color = MSO_THEME_COLOR.DARK_1
    major_paragraph = shape.text_frame.paragraphs[0]
    major_paragraph.text = "Theme major font"
    major_run = major_paragraph.runs[0]
    major_run.font.name = "+mj-lt"
    major_run.font.size = Pt(28)
    minor_paragraph = shape.text_frame.add_paragraph()
    minor_paragraph.text = "Theme minor font"
    minor_run = minor_paragraph.runs[0]
    minor_run.font.name = "+mn-lt"
    minor_run.font.size = Pt(28)

    path = os.path.join(OUTPUT_DIR, filename)
    prs.save(path)
    major_typeface = "Carlito" if expected else "Liberation Sans"
    minor_typeface = "Caladea" if expected else "Liberation Serif"
    replace_zip_part_text(
        path,
        "ppt/theme/theme1.xml",
        '<a:majorFont><a:latin typeface="Calibri"/>',
        f'<a:majorFont><a:latin typeface="{major_typeface}"/>',
    )
    replace_zip_part_text(
        path,
        "ppt/theme/theme1.xml",
        '<a:minorFont><a:latin typeface="Calibri"/>',
        f'<a:minorFont><a:latin typeface="{minor_typeface}"/>',
    )
    if expected:
        replace_zip_part_text(
            path,
            "ppt/theme/theme1.xml",
            'val="4F81BD"',
            'val="C00000"',
        )
    print(f"  Created: {filename}")


def replace_zip_part_text(path, part_path, old, new):
    """Replace one UTF-8 fragment in a ZIP part without creating duplicate entries."""
    with zipfile.ZipFile(path, "r") as source:
        entries = [(item, source.read(item.filename)) for item in source.infolist()]
    replaced = False
    with tempfile.NamedTemporaryFile(
        suffix=".pptx", dir=os.path.dirname(path), delete=False
    ) as tmp:
        temporary_path = tmp.name
    try:
        with zipfile.ZipFile(temporary_path, "w") as output:
            for item, data in entries:
                if item.filename == part_path:
                    text = data.decode("utf-8")
                    if old not in text:
                        raise RuntimeError(f"theme fixture fragment not found: {old}")
                    data = text.replace(old, new, 1).encode("utf-8")
                    replaced = True
                output.writestr(item, data)
        if not replaced:
            raise RuntimeError(f"theme fixture part not found: {part_path}")
        os.chmod(temporary_path, 0o644)
        os.replace(temporary_path, path)
    finally:
        if os.path.exists(temporary_path):
            os.unlink(temporary_path)


def create_editor_validity_fixtures():
    """PPTX source / expected pairs consumed by editor-validity.test.ts."""
    create_editor_validity_text_fixture(
        "editor-validity-text-source.pptx",
        "Original LibreOffice text",
    )
    create_editor_validity_theme_fixture(
        "editor-validity-theme-source.pptx",
        expected=False,
    )
    create_editor_validity_theme_fixture(
        "editor-validity-theme-expected.pptx",
        expected=True,
    )
    create_editor_validity_text_fixture(
        "editor-validity-text-expected.pptx",
        "Edited LibreOffice text",
    )
    create_editor_validity_transform_fixture(
        "editor-validity-transform-source.pptx",
        Inches(0.9),
        Inches(1.2),
        Inches(2.4),
        Inches(1.2),
    )
    create_editor_validity_transform_fixture(
        "editor-validity-transform-expected.pptx",
        Inches(3.0),
        Inches(2.1),
        Inches(3.2),
        Inches(1.6),
    )
    create_editor_validity_formatting_fixture(
        "editor-validity-formatting-source.pptx",
        expected=False,
    )
    create_editor_validity_formatting_fixture(
        "editor-validity-formatting-expected.pptx",
        expected=True,
    )
    create_editor_validity_paragraph_fixture(
        "editor-validity-paragraph-source.pptx",
        expected=False,
    )
    create_editor_validity_paragraph_fixture(
        "editor-validity-paragraph-expected.pptx",
        expected=True,
    )
    create_editor_validity_image_fixture(
        "editor-validity-image-source.pptx",
        "iVBORw0KGgoAAAANSUhEUgAAAAQAAAAECAIAAAAmkwkpAAAAEUlEQVR4nGP8z4AATEhsPBwAM9EBBzDn4UwAAAAASUVORK5CYII=",
    )
    create_editor_validity_image_fixture(
        "editor-validity-image-expected.pptx",
        "iVBORw0KGgoAAAANSUhEUgAAAAQAAAAECAIAAAAmkwkpAAAAE0lEQVR4nGNkYPjPAANMcBZeDgAx0wEH1s7nlgAAAABJRU5ErkJggg==",
    )
    create_editor_validity_picture_crop_fixture(
        "editor-validity-picture-crop-source.pptx",
        expected=False,
    )
    create_editor_validity_picture_crop_fixture(
        "editor-validity-picture-crop-expected.pptx",
        expected=True,
    )
    create_editor_validity_chart_fixture(
        "editor-validity-chart-source.pptx",
        expected=False,
    )
    create_editor_validity_chart_fixture(
        "editor-validity-chart-expected.pptx",
        expected=True,
    )
    create_editor_validity_chart_removal_fixture(
        "editor-validity-chart-removal-source.pptx",
        expected=False,
    )
    create_editor_validity_category_combo_chart_fixture(
        "editor-validity-category-combo-chart-source.pptx",
        expected=False,
    )
    create_editor_validity_category_combo_chart_fixture(
        "editor-validity-category-combo-chart-expected.pptx",
        expected=True,
    )
    create_editor_validity_chart_removal_fixture(
        "editor-validity-chart-removal-expected.pptx",
        expected=True,
    )
    create_editor_validity_scatter_chart_fixture(
        "editor-validity-scatter-chart-source.pptx",
        expected=False,
    )
    create_editor_validity_scatter_chart_fixture(
        "editor-validity-scatter-chart-expected.pptx",
        expected=True,
    )
    create_editor_validity_bubble_chart_fixture(
        "editor-validity-bubble-chart-source.pptx",
        expected=False,
    )
    create_editor_validity_bubble_chart_fixture(
        "editor-validity-bubble-chart-expected.pptx",
        expected=True,
    )
    create_editor_validity_table_text_fixture(
        "editor-validity-table-text-source.pptx",
        "Original LibreOffice table text",
    )
    create_editor_validity_table_text_fixture(
        "editor-validity-table-text-expected.pptx",
        "Edited LibreOffice table text",
    )
    create_editor_validity_table_cell_properties_fixture(
        "editor-validity-table-cell-properties-source.pptx",
        expected=False,
    )
    create_editor_validity_table_cell_properties_fixture(
        "editor-validity-table-cell-properties-expected.pptx",
        expected=True,
    )
    create_editor_validity_group_fixture(
        "editor-validity-group-source.pptx",
        grouped=False,
    )
    create_editor_validity_group_fixture(
        "editor-validity-group-expected.pptx",
        grouped=True,
    )
    create_editor_validity_affine_move_fixture()
    create_editor_validity_cross_slide_move_fixture(
        "editor-validity-cross-slide-move-source.pptx",
        expected=False,
    )
    create_editor_validity_cross_slide_move_fixture(
        "editor-validity-cross-slide-move-expected.pptx",
        expected=True,
    )
    create_editor_validity_drawing_delete_fixture(
        "editor-validity-drawing-delete-source.pptx",
    )
    create_editor_validity_group_fixture(
        "editor-validity-nested-group-delete-source.pptx",
        grouped=True,
    )


def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    print("Generating editor-validity fixtures...")
    create_basic_shapes()
    create_editor_validity_fixtures()
    print("Done!")


if __name__ == "__main__":
    main()

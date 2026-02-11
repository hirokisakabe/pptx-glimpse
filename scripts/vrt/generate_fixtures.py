#!/usr/bin/env python3
"""
LibreOffice VRT 用の PPTX フィクスチャを python-pptx で生成する。

Usage:
    python3 scripts/vrt/generate_fixtures.py
"""

import os

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


def make_element(tag, **attribs):
    """OOXML 要素を作成する (e.g., make_element("a:gradFill"))"""
    element = etree.SubElement(etree.Element("dummy"), qn(tag))
    element.getparent().remove(element)
    for key, val in attribs.items():
        element.set(key, val)
    return element

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "../../tests/vrt/libreoffice-fixtures")

SLIDE_WIDTH = 9144000
SLIDE_HEIGHT = 5143500


def new_presentation():
    """スライドサイズを固定した新しいプレゼンテーションを生成する。"""
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH)
    prs.slide_height = Emu(SLIDE_HEIGHT)
    return prs


def create_basic_shapes():
    """基本図形: rect, ellipse, roundRect (ソリッド塗り + 枠線)"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    shapes_def = [
        (MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(2.5), Inches(2),
         RGBColor(0x44, 0x72, 0xC4), "Rectangle"),
        (MSO_SHAPE.OVAL, Inches(3.5), Inches(0.5), Inches(2.5), Inches(2),
         RGBColor(0xED, 0x7D, 0x31), "Oval"),
        (MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(0.5), Inches(2.5), Inches(2),
         RGBColor(0xA5, 0xA5, 0xA5), "Rounded Rect"),
        (MSO_SHAPE.DIAMOND, Inches(0.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0xFF, 0xC0, 0x00), "Diamond"),
        (MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0x5B, 0x9B, 0xD5), "Triangle"),
        (MSO_SHAPE.HEXAGON, Inches(6.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0x70, 0xAD, 0x47), "Hexagon"),
    ]

    for shape_type, left, top, width, height, color, _label in shapes_def:
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    prs.save(os.path.join(OUTPUT_DIR, "lo-basic-shapes.pptx"))
    print("  Created: lo-basic-shapes.pptx")


def create_text_formatting():
    """テキスト書式: 太字、イタリック、フォントサイズ、色、配置"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tests = [
        {"text": "Bold Text", "bold": True, "size": Pt(24)},
        {"text": "Italic Text", "italic": True, "size": Pt(24)},
        {"text": "Large 36pt", "size": Pt(36)},
        {"text": "Red Text", "size": Pt(24), "color": RGBColor(0xFF, 0x00, 0x00)},
        {"text": "Center Aligned", "size": Pt(24), "align": PP_ALIGN.CENTER},
        {"text": "Right Aligned", "size": Pt(24), "align": PP_ALIGN.RIGHT},
    ]

    for i, t in enumerate(tests):
        col = i % 2
        row = i // 2
        left = Inches(0.3 + col * 4.7)
        top = Inches(0.3 + row * 1.7)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(4.4), Inches(1.4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        shape.line.width = Pt(0.5)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t["text"]

        if "align" in t:
            p.alignment = t["align"]

        run = p.runs[0]
        run.font.name = "Liberation Sans"
        run.font.size = t.get("size", Pt(18))

        if t.get("bold"):
            run.font.bold = True
        if t.get("italic"):
            run.font.italic = True
        if "color" in t:
            run.font.color.rgb = t["color"]

    prs.save(os.path.join(OUTPUT_DIR, "lo-text-formatting.pptx"))
    print("  Created: lo-text-formatting.pptx")


def create_fill_and_lines():
    """塗りと線: ソリッド塗り各色 + 枠線スタイル"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    colors = [
        (RGBColor(0xFF, 0x63, 0x84), "Pink"),
        (RGBColor(0x36, 0xA2, 0xEB), "Blue"),
        (RGBColor(0xFF, 0xCE, 0x56), "Yellow"),
    ]

    for i, (color, _label) in enumerate(colors):
        left = Inches(0.5 + i * 3)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(0.5),
            Inches(2.5), Inches(2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(2)

    # 枠線のみの図形（塗りなし）
    no_fill_shapes = [
        (MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0x44, 0x72, 0xC4), Pt(1)),
        (MSO_SHAPE.OVAL, Inches(3.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0xED, 0x7D, 0x31), Pt(2)),
        (MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(3), Inches(2.5), Inches(2),
         RGBColor(0x70, 0xAD, 0x47), Pt(3)),
    ]

    for shape_type, left, top, width, height, line_color, line_width in no_fill_shapes:
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.background()
        shape.line.color.rgb = line_color
        shape.line.width = line_width

    prs.save(os.path.join(OUTPUT_DIR, "lo-fill-and-lines.pptx"))
    print("  Created: lo-fill-and-lines.pptx")


def create_gradient_fills():
    """グラデーション塗り: 水平/垂直/対角/3色/枠線付き/暗色"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    gradients = [
        {
            "shape": MSO_SHAPE.ROUNDED_RECTANGLE,
            "stops": [(0, "4472C4"), (100000, "ED7D31")],
            "angle": 0,
            "label": "Horizontal",
        },
        {
            "shape": MSO_SHAPE.ROUNDED_RECTANGLE,
            "stops": [(0, "FF0000"), (100000, "FFFF00")],
            "angle": 5400000,
            "label": "Vertical",
        },
        {
            "shape": MSO_SHAPE.ROUNDED_RECTANGLE,
            "stops": [(0, "70AD47"), (100000, "7030A0")],
            "angle": 2700000,
            "label": "Diagonal",
        },
        {
            "shape": MSO_SHAPE.RECTANGLE,
            "stops": [(0, "FF0000"), (50000, "FFFFFF"), (100000, "4472C4")],
            "angle": 0,
            "label": "3-Color",
        },
        {
            "shape": MSO_SHAPE.OVAL,
            "stops": [(0, "5B9BD5"), (100000, "FFC000")],
            "angle": 0,
            "label": "With Border",
            "border": True,
        },
        {
            "shape": MSO_SHAPE.DIAMOND,
            "stops": [(0, "000000"), (100000, "808080")],
            "angle": 5400000,
            "label": "Dark",
        },
    ]

    for i, g in enumerate(gradients):
        col = i % 3
        row = i // 3
        left = Inches(0.3 + col * 3.2)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(
            g["shape"], left, top, Inches(2.8), Inches(2.2)
        )
        shape.fill.background()

        # XML でグラデーション塗りを設定
        spPr = shape._element.spPr
        # noFill を削除
        for child in list(spPr):
            tag = etree.QName(child.tag).localname
            if tag in ("noFill", "solidFill"):
                spPr.remove(child)

        gradFill = make_element("a:gradFill")
        gsLst = make_element("a:gsLst")
        for pos, color_hex in g["stops"]:
            gs = make_element("a:gs")
            gs.set("pos", str(pos))
            srgbClr = make_element("a:srgbClr")
            srgbClr.set("val", color_hex)
            gs.append(srgbClr)
            gsLst.append(gs)
        gradFill.append(gsLst)

        lin = make_element("a:lin")
        lin.set("ang", str(g["angle"]))
        lin.set("scaled", "1")
        gradFill.append(lin)
        spPr.append(gradFill)

        if g.get("border"):
            shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
            shape.line.width = Pt(2)

    prs.save(os.path.join(OUTPUT_DIR, "lo-gradient-fills.pptx"))
    print("  Created: lo-gradient-fills.pptx")


def create_dash_lines():
    """線のダッシュスタイル: solid, dash, dot, dashDot, lgDash, lgDashDot"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    dash_styles = [
        ("solid", "Solid"),
        ("dash", "Dash"),
        ("dot", "Dot"),
        ("dashDot", "DashDot"),
        ("lgDash", "LgDash"),
        ("lgDashDot", "LgDashDot"),
    ]

    for i, (dash_val, label) in enumerate(dash_styles):
        col = i % 3
        row = i // 3
        left = Inches(0.3 + col * 3.2)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(2.8), Inches(2.2)
        )
        shape.fill.background()
        shape.line.color.rgb = RGBColor(0x44, 0x72, 0xC4)
        shape.line.width = Pt(2.5)

        # XML でダッシュスタイルを設定
        ln = shape.line._ln
        prstDash = make_element("a:prstDash")
        prstDash.set("val", dash_val)
        ln.append(prstDash)

        # ラベルテキスト
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Liberation Sans"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(os.path.join(OUTPUT_DIR, "lo-dash-lines.pptx"))
    print("  Created: lo-dash-lines.pptx")


def create_text_decoration():
    """テキスト装飾: 下線、取り消し線、太字+下線、イタリック+取り消し線、複数行、上付き"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tests = [
        {"text": "Underline Text", "underline": True},
        {"text": "Strikethrough", "strikethrough": True},
        {"text": "Bold + Underline", "bold": True, "underline": True},
        {"text": "Italic + Strike", "italic": True, "strikethrough": True},
        {"text": "Line 1\nLine 2\nLine 3", "multiline": True},
        {"text": "Normal", "superscript": "sup", "sup_text": "super"},
    ]

    for i, t in enumerate(tests):
        col = i % 3
        row = i // 3
        left = Inches(0.3 + col * 3.2)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(2.8), Inches(2.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
        shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        shape.line.width = Pt(0.5)

        tf = shape.text_frame
        tf.word_wrap = True

        if t.get("multiline"):
            lines = t["text"].split("\n")
            for j, line_text in enumerate(lines):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line_text
                run = p.runs[0]
                run.font.name = "Liberation Sans"
                run.font.size = Pt(18)
        elif t.get("superscript"):
            p = tf.paragraphs[0]
            p.text = t["text"]
            run = p.runs[0]
            run.font.name = "Liberation Sans"
            run.font.size = Pt(20)
            # 上付き文字を追加
            sup_run = p.add_run()
            sup_run.text = t["sup_text"]
            sup_run.font.name = "Liberation Sans"
            sup_run.font.size = Pt(20)
            rPr = sup_run._r.get_or_add_rPr()
            rPr.set("baseline", "30000")
        else:
            p = tf.paragraphs[0]
            p.text = t["text"]
            run = p.runs[0]
            run.font.name = "Liberation Sans"
            run.font.size = Pt(20)

            if t.get("bold"):
                run.font.bold = True
            if t.get("italic"):
                run.font.italic = True
            if t.get("underline"):
                run.font.underline = True
            if t.get("strikethrough"):
                rPr = run._r.get_or_add_rPr()
                rPr.set("strike", "sngStrike")

    prs.save(os.path.join(OUTPUT_DIR, "lo-text-decoration.pptx"))
    print("  Created: lo-text-decoration.pptx")


def create_tables():
    """テーブル: ヘッダ行 + データ行 (交互背景色)"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    rows, cols = 4, 3
    left = Inches(1)
    top = Inches(0.8)
    width = Inches(8)
    height = Inches(3.5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    headers = ["Name", "Category", "Value"]
    data = [
        ["Alpha", "Type A", "100"],
        ["Beta", "Type B", "200"],
        ["Gamma", "Type C", "300"],
    ]

    # ヘッダ行
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Liberation Sans"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # データ行
    for i, row_data in enumerate(data):
        bg = RGBColor(0xF2, 0xF2, 0xF2) if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        for j, value in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0]
            run.font.name = "Liberation Sans"
            run.font.size = Pt(14)

    prs.save(os.path.join(OUTPUT_DIR, "lo-tables.pptx"))
    print("  Created: lo-tables.pptx")


def create_bullets():
    """箇条書き・番号付きリスト: buChar (丸/ダッシュ) + buAutoNum (数字/アルファベット)"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bullet_configs = [
        {
            "title": "Bullet (dot)",
            "type": "buChar",
            "char": "\u2022",
            "items": ["First item", "Second item", "Third item"],
        },
        {
            "title": "Bullet (dash)",
            "type": "buChar",
            "char": "-",
            "items": ["Apple", "Banana", "Cherry"],
        },
        {
            "title": "Numbered (1. 2. 3.)",
            "type": "buAutoNum",
            "scheme": "arabicPeriod",
            "items": ["Step one", "Step two", "Step three"],
        },
        {
            "title": "Alpha (a. b. c.)",
            "type": "buAutoNum",
            "scheme": "alphaLcPeriod",
            "items": ["Option A", "Option B", "Option C"],
        },
    ]

    for i, config in enumerate(bullet_configs):
        col = i % 2
        row = i // 2
        left = Inches(0.3 + col * 4.7)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(4.4), Inches(2.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        shape.line.width = Pt(0.5)

        tf = shape.text_frame
        tf.word_wrap = True

        # タイトル行
        p = tf.paragraphs[0]
        p.text = config["title"]
        run = p.runs[0]
        run.font.name = "Liberation Sans"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)

        # 箇条書き項目
        for item_text in config["items"]:
            p = tf.add_paragraph()
            p.text = item_text
            p.level = 0
            run = p.runs[0]
            run.font.name = "Liberation Sans"
            run.font.size = Pt(14)

            pPr = p._pPr
            if pPr is None:
                pPr = make_element("a:pPr")
                p._p.insert(0, pPr)
            pPr.set("marL", str(Inches(0.3)))
            pPr.set("indent", str(-Inches(0.2)))

            if config["type"] == "buChar":
                buChar = make_element("a:buChar")
                buChar.set("char", config["char"])
                pPr.append(buChar)
            elif config["type"] == "buAutoNum":
                buAutoNum = make_element("a:buAutoNum")
                buAutoNum.set("type", config["scheme"])
                pPr.append(buAutoNum)

    prs.save(os.path.join(OUTPUT_DIR, "lo-bullets.pptx"))
    print("  Created: lo-bullets.pptx")


def create_transforms():
    """回転・フリップ: 45度/90度回転、水平/垂直フリップ、組み合わせ"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    transforms = [
        {"rotation": 45.0, "flipH": False, "flipV": False, "label": "Rot 45"},
        {"rotation": 90.0, "flipH": False, "flipV": False, "label": "Rot 90"},
        {"rotation": 0.0, "flipH": True, "flipV": False, "label": "FlipH"},
        {"rotation": 0.0, "flipH": False, "flipV": True, "label": "FlipV"},
        {"rotation": 0.0, "flipH": True, "flipV": True, "label": "FlipHV"},
        {"rotation": 30.0, "flipH": True, "flipV": False, "label": "Rot30+FlipH"},
    ]

    for i, t in enumerate(transforms):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 3.2)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, left, top, Inches(2.5), Inches(1.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1)

        if t["rotation"] != 0.0:
            shape.rotation = t["rotation"]

        xfrm = shape._element.spPr.xfrm
        if t["flipH"]:
            xfrm.set("flipH", "1")
        if t["flipV"]:
            xfrm.set("flipV", "1")

    prs.save(os.path.join(OUTPUT_DIR, "lo-transforms.pptx"))
    print("  Created: lo-transforms.pptx")


def create_groups():
    """グループ図形: 2グループ (各2図形)"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # python-pptx でのグループ作成は XML 操作が必要
    # グループ1: 青矩形 + オレンジ楕円
    group_shapes = [
        {
            "group_left": Inches(0.5),
            "group_top": Inches(0.8),
            "group_w": Inches(4),
            "group_h": Inches(3.5),
            "children": [
                {
                    "type": MSO_SHAPE.RECTANGLE,
                    "left": Inches(0.5),
                    "top": Inches(0.8),
                    "w": Inches(1.8),
                    "h": Inches(1.5),
                    "color": RGBColor(0x44, 0x72, 0xC4),
                },
                {
                    "type": MSO_SHAPE.OVAL,
                    "left": Inches(2.5),
                    "top": Inches(2.5),
                    "w": Inches(1.8),
                    "h": Inches(1.5),
                    "color": RGBColor(0xED, 0x7D, 0x31),
                },
            ],
        },
        {
            "group_left": Inches(5.2),
            "group_top": Inches(0.8),
            "group_w": Inches(4),
            "group_h": Inches(3.5),
            "children": [
                {
                    "type": MSO_SHAPE.DIAMOND,
                    "left": Inches(5.2),
                    "top": Inches(0.8),
                    "w": Inches(1.8),
                    "h": Inches(1.5),
                    "color": RGBColor(0x70, 0xAD, 0x47),
                },
                {
                    "type": MSO_SHAPE.ISOSCELES_TRIANGLE,
                    "left": Inches(7.2),
                    "top": Inches(2.5),
                    "w": Inches(1.8),
                    "h": Inches(1.5),
                    "color": RGBColor(0xFF, 0xC0, 0x00),
                },
            ],
        },
    ]

    for group_def in group_shapes:
        # グループ図形を XML で構築
        spTree = slide.shapes._spTree

        grpSp = make_element("p:grpSp")

        # grpSpPr (グループの位置・サイズ)
        grpSpPr = make_element("p:grpSpPr")
        xfrm = make_element("a:xfrm")
        off = make_element("a:off")
        off.set("x", str(int(group_def["group_left"])))
        off.set("y", str(int(group_def["group_top"])))
        xfrm.append(off)
        ext = make_element("a:ext")
        ext.set("cx", str(int(group_def["group_w"])))
        ext.set("cy", str(int(group_def["group_h"])))
        xfrm.append(ext)
        chOff = make_element("a:chOff")
        chOff.set("x", str(int(group_def["group_left"])))
        chOff.set("y", str(int(group_def["group_top"])))
        xfrm.append(chOff)
        chExt = make_element("a:chExt")
        chExt.set("cx", str(int(group_def["group_w"])))
        chExt.set("cy", str(int(group_def["group_h"])))
        xfrm.append(chExt)
        grpSpPr.append(xfrm)
        grpSp.append(grpSpPr)

        # nvGrpSpPr
        nvGrpSpPr = make_element("p:nvGrpSpPr")
        cNvPr = make_element("p:cNvPr")
        cNvPr.set("id", str(100 + group_shapes.index(group_def)))
        cNvPr.set("name", f"Group {group_shapes.index(group_def) + 1}")
        nvGrpSpPr.append(cNvPr)
        cNvGrpSpPr = make_element("p:cNvGrpSpPr")
        nvGrpSpPr.append(cNvGrpSpPr)
        nvPr = make_element("p:nvPr")
        nvGrpSpPr.append(nvPr)
        grpSp.insert(0, nvGrpSpPr)

        # 子図形を追加
        for ci, child in enumerate(group_def["children"]):
            sp = make_element("p:sp")

            # nvSpPr
            nvSpPr = make_element("p:nvSpPr")
            cNvPr2 = make_element("p:cNvPr")
            cNvPr2.set("id", str(200 + group_shapes.index(group_def) * 10 + ci))
            cNvPr2.set("name", f"Shape {ci + 1}")
            nvSpPr.append(cNvPr2)
            cNvSpPr = make_element("p:cNvSpPr")
            nvSpPr.append(cNvSpPr)
            nvPr2 = make_element("p:nvPr")
            nvSpPr.append(nvPr2)
            sp.append(nvSpPr)

            # spPr
            spPr = make_element("p:spPr")
            childXfrm = make_element("a:xfrm")
            childOff = make_element("a:off")
            childOff.set("x", str(int(child["left"])))
            childOff.set("y", str(int(child["top"])))
            childXfrm.append(childOff)
            childExt = make_element("a:ext")
            childExt.set("cx", str(int(child["w"])))
            childExt.set("cy", str(int(child["h"])))
            childXfrm.append(childExt)
            spPr.append(childXfrm)

            # prstGeom (プリセット図形)
            prstGeom = make_element("a:prstGeom")
            shape_map = {
                MSO_SHAPE.RECTANGLE: "rect",
                MSO_SHAPE.OVAL: "ellipse",
                MSO_SHAPE.DIAMOND: "diamond",
                MSO_SHAPE.ISOSCELES_TRIANGLE: "triangle",
            }
            prstGeom.set("prst", shape_map[child["type"]])
            avLst = make_element("a:avLst")
            prstGeom.append(avLst)
            spPr.append(prstGeom)

            # solidFill
            solidFill = make_element("a:solidFill")
            srgbClr = make_element("a:srgbClr")
            color = child["color"]
            srgbClr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
            solidFill.append(srgbClr)
            spPr.append(solidFill)

            # ln (枠線)
            ln = make_element("a:ln")
            ln.set("w", str(int(Pt(1.5))))
            lnFill = make_element("a:solidFill")
            lnClr = make_element("a:srgbClr")
            lnClr.set("val", "333333")
            lnFill.append(lnClr)
            ln.append(lnFill)
            spPr.append(ln)

            sp.append(spPr)
            grpSp.append(sp)

        spTree.append(grpSp)

    prs.save(os.path.join(OUTPUT_DIR, "lo-groups.pptx"))
    print("  Created: lo-groups.pptx")


def create_slide_background():
    """スライド背景: ソリッドカラー背景 + 白テキストボックス + 枠線のみ図形"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # スライド背景色を設定
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)

    # 白い矩形 + テキスト
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(3.5), Inches(1.5)
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape1.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape1.line.width = Pt(0)
    tf = shape1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "White Box"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Liberation Sans"
    run.font.size = Pt(24)
    run.font.bold = True

    # 黄色い矩形
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(1), Inches(3.5), Inches(1.5)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0xFF, 0xCE, 0x56)
    shape2.line.width = Pt(0)
    tf2 = shape2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Yellow Box"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Liberation Sans"
    run2.font.size = Pt(24)
    run2.font.bold = True

    # 枠線のみの図形 (白枠)
    shape3 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(2.5), Inches(3), Inches(5), Inches(1.8)
    )
    shape3.fill.background()
    shape3.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape3.line.width = Pt(3)

    prs.save(os.path.join(OUTPUT_DIR, "lo-slide-background.pptx"))
    print("  Created: lo-slide-background.pptx")


def create_flowchart_shapes():
    """フローチャート図形: 8つのフローチャートプリセット (ソリッド塗り + 枠線)"""
    prs = new_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    flowchart_defs = [
        (MSO_SHAPE.FLOWCHART_PROCESS, RGBColor(0x44, 0x72, 0xC4), "Process"),
        (MSO_SHAPE.FLOWCHART_ALTERNATE_PROCESS, RGBColor(0xED, 0x7D, 0x31), "Alt Process"),
        (MSO_SHAPE.FLOWCHART_DECISION, RGBColor(0xA5, 0xA5, 0xA5), "Decision"),
        (MSO_SHAPE.FLOWCHART_DATA, RGBColor(0xFF, 0xC0, 0x00), "Data"),
        (MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS, RGBColor(0x5B, 0x9B, 0xD5), "Predefined"),
        (MSO_SHAPE.FLOWCHART_DOCUMENT, RGBColor(0x70, 0xAD, 0x47), "Document"),
        (MSO_SHAPE.FLOWCHART_TERMINATOR, RGBColor(0xFF, 0x63, 0x84), "Terminator"),
        (MSO_SHAPE.FLOWCHART_PREPARATION, RGBColor(0x36, 0xA2, 0xEB), "Preparation"),
    ]

    for i, (shape_type, color, _label) in enumerate(flowchart_defs):
        col = i % 4
        row = i // 4
        left = Inches(0.3 + col * 2.4)
        top = Inches(0.3 + row * 2.6)

        shape = slide.shapes.add_shape(shape_type, left, top, Inches(2.1), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

    prs.save(os.path.join(OUTPUT_DIR, "lo-flowchart-shapes.pptx"))
    print("  Created: lo-flowchart-shapes.pptx")


def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    print("Generating LibreOffice VRT fixtures...")
    create_basic_shapes()
    create_text_formatting()
    create_fill_and_lines()
    create_gradient_fills()
    create_dash_lines()
    create_text_decoration()
    create_tables()
    create_bullets()
    create_transforms()
    create_groups()
    create_slide_background()
    create_flowchart_shapes()
    print("Done!")


if __name__ == "__main__":
    main()
